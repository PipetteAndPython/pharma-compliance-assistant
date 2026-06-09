# =========================
# DOCUMENT LOADER
# =========================
# Supports: PDF (.pdf), Word (.docx), PowerPoint (.pptx)

from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from pathlib import Path


# =========================
# FORMAT-SPECIFIC LOADERS
# =========================

def load_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def load_docx(file_path):
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text


def load_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


# =========================
# SUPPORTED EXTENSIONS
# =========================

SUPPORTED_EXTENSIONS = {
    ".pdf":  load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
}


# =========================
# MAIN LOADER
# =========================

def load_documents(folder_path):
    """
    Loads all supported documents from a folder.
    Returns a list of dicts with 'file' and 'text' keys.
    """

    folder = Path(folder_path)
    docs = []
    skipped = []

    for file in sorted(folder.iterdir()):

        ext = file.suffix.lower()

        if ext not in SUPPORTED_EXTENSIONS:
            skipped.append(file.name)
            continue

        try:
            loader = SUPPORTED_EXTENSIONS[ext]
            text = loader(file)

            if not text.strip():
                print(f"⚠ No extractable text: {file.name}")
                continue

            docs.append({
                "file": file.name,
                "text": text
            })

        except Exception as e:
            print(f"❌ Error reading {file.name}: {e}")

    if skipped:
        print(f"⚠ Skipped files (unsupported format): {skipped}")

    print(f"✔ Documents loaded successfully: {len(docs)}")

    return docs